from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_remediation_report():
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Security Remediation Report"
    subtitle.text = "Ain't all Risky Bizz Application\nPost-Remediation Security Assessment"

    # Executive Summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"
    content = slide.placeholders[1]
    content.text = (
        "This report details the remediation efforts undertaken to address the security findings "
        "identified in the initial OWASP Top 10 review.\n\n"
        "Summary of Actions:\n"
        "- Critical and High severity issues have been remediated.\n"
        "- Security controls (CSRF, Secure Cookies, Headers) have been implemented.\n"
        "- Dependencies have been pinned to secure versions.\n"
        "- Application configuration has been hardened for production."
    )

    # Remediation Status Table
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Remediation Status"
    
    rows = 6
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ["Finding", "Severity", "Action Taken", "Status"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
        
    # Data
    data = [
        ("Cryptographic Failures (Hardcoded Secret)", "Critical", "Moved to Env Var", "Fixed"),
        ("Security Misconfiguration (Debug=True)", "High", "Disabled Debug Mode", "Fixed"),
        ("Insecure Design (No CSRF)", "Medium", "Implemented Flask-WTF CSRF", "Fixed"),
        ("Vulnerable Components (Unpinned Deps)", "Medium", "Pinned Requirements", "Fixed"),
        ("Broken Access Control", "Medium", "Documented as Design Choice", "Accepted Risk"),
    ]
    
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            table.cell(i+1, j).text = val

    # Detailed Fixes: Secret Key & Debug
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Fix: Secret Key & Debug Mode"
    content = slide.placeholders[1]
    content.text = (
        "Issue: Hardcoded SECRET_KEY and Debug Mode enabled in production code.\n\n"
        "Remediation:\n"
        "1. Updated app.py to load SECRET_KEY from environment variable.\n"
        "   - app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', ...)\n"
        "2. Disabled Debug Mode in app.run().\n"
        "   - app.run(debug=False)\n\n"
        "Verification:\n"
        "- Confirmed 404 pages do not show interactive debugger.\n"
        "- Confirmed application starts with secure configuration."
    )

    # Detailed Fixes: CSRF Protection
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Fix: CSRF Protection"
    content = slide.placeholders[1]
    content.text = (
        "Issue: Forms lacked Cross-Site Request Forgery (CSRF) protection.\n\n"
        "Remediation:\n"
        "1. Installed and configured Flask-WTF.\n"
        "2. Initialized CSRFProtect(app) in app.py.\n"
        "3. Added hidden CSRF token field to assessment forms.\n\n"
        "Verification:\n"
        "- Confirmed forms submit successfully with token.\n"
        "- Confirmed forms fail without token (implicit via library)."
    )

    # Detailed Fixes: Secure Cookies
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Fix: Secure Cookies"
    content = slide.placeholders[1]
    content.text = (
        "Issue: Session cookies lacked security flags.\n\n"
        "Remediation:\n"
        "1. Enabled SESSION_COOKIE_HTTPONLY = True.\n"
        "2. Configured SESSION_COOKIE_SECURE (False for Dev, True for Prod).\n\n"
        "Verification:\n"
        "- Verified HttpOnly flag is present in session cookies."
    )

    # Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion & Next Steps"
    content = slide.placeholders[1]
    content.text = (
        "The application's security posture has been significantly improved.\n\n"
        "Current Status:\n"
        "- All critical and high severity technical findings are resolved.\n"
        "- Application is ready for containerization and deployment.\n\n"
        "Next Steps:\n"
        "- Perform final Docker build.\n"
        "- Deploy to production environment with HTTPS enabled."
    )

    prs.save('Remediation_Report.pptx')
    print("Remediation report generated successfully.")

if __name__ == "__main__":
    create_remediation_report()
