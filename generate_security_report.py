from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper to set title and content
    def add_slide(title_text, content_text_list, level=0):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear() # Clear existing paragraphs
        
        for text in content_text_list:
            p = tf.add_paragraph()
            p.text = text
            p.level = level

    # Helper for findings slides with severity color
    def add_finding_slide(title, severity, description, evidence, recommendation):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = f"{title} ({severity})"
        
        # Color code severity in title if possible, or just text
        # Simple text frame population
        tf = slide.shapes.placeholders[1].text_frame
        
        p = tf.add_paragraph()
        p.text = "Description:"
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = description
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Evidence:"
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = evidence
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Recommendation:"
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = recommendation
        p.level = 1

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "Security Review Report"
    subtitle.text = "Ain't all Risky Bizz Application\nOWASP Top 10 Assessment\nNovember 28, 2025"

    # 2. Executive Summary
    add_slide("Executive Summary", [
        "A comprehensive security review was performed based on the OWASP Top 10 (2021) standard.",
        "Key Findings:",
        "• 1 Critical Issue (Cryptographic Failures)",
        "• 1 High Issue (Security Misconfiguration)",
        "• 3 Medium Issues (Insecure Design, Components, Access Control)",
        "• Positive Results: No XSS or SQL Injection vulnerabilities found."
    ])

    # 3. Finding: Cryptographic Failures
    add_finding_slide(
        "A02: Cryptographic Failures",
        "CRITICAL",
        "The application uses a hardcoded SECRET_KEY ('dev-secret-key') in the source code. This compromises all session security, allowing attackers to forge session cookies.",
        "app.py line 6: app.config['SECRET_KEY'] = 'dev-secret-key'",
        "Use an environment variable for the secret key. Ensure it is a long, random string. Never commit secrets to version control."
    )

    # 4. Finding: Security Misconfiguration
    add_finding_slide(
        "A05: Security Misconfiguration",
        "HIGH",
        "Debug mode is enabled in the application configuration. This can leak sensitive stack traces, environment variables, and code snippets to attackers if an error occurs.",
        "app.py line 66: app.run(debug=True)",
        "Disable debug mode in production. Use app.run(debug=False) or rely on a production WSGI server configuration."
    )

    # 5. Finding: Insecure Design (CSRF)
    add_finding_slide(
        "A04: Insecure Design",
        "MEDIUM",
        "The application accepts POST requests (e.g., in the assessment wizard) without Anti-CSRF tokens. Attackers could trick users into submitting unauthorized data.",
        "HTML forms in assessment_wizard.html lack {{ csrf_token() }}.",
        "Implement Flask-WTF or a similar library to generate and validate CSRF tokens for all state-changing forms."
    )

    # 6. Finding: Vulnerable Components
    add_finding_slide(
        "A06: Vulnerable Components",
        "MEDIUM",
        "Dependency versions are not pinned in requirements.txt (e.g., 'flask' instead of 'flask==2.3.2'). This leads to non-reproducible builds and potential use of vulnerable versions.",
        "requirements.txt contains unpinned package names.",
        "Pin exact versions in requirements.txt (e.g., Flask==3.0.0) and regularly scan dependencies for vulnerabilities."
    )

    # 7. Finding: Broken Access Control
    add_finding_slide(
        "A01: Broken Access Control",
        "MEDIUM",
        "No authentication or authorization mechanisms are implemented. The application is open to any user with network access.",
        "No @login_required decorators or user management logic in app.py.",
        "Implement a user authentication system if the application is intended for restricted use. Ensure proper session management."
    )

    # 8. Positive Findings
    add_slide("Positive Findings", [
        "• A03: Injection (Pass)",
        "  - Reflected XSS testing was unsuccessful (Jinja2 auto-escaping active).",
        "  - SQL Injection risk is low due to SQLAlchemy ORM usage.",
        "• A10: SSRF (Pass)",
        "  - No functionality found that fetches remote resources based on user input."
    ])

    # 9. Next Steps
    add_slide("Recommended Next Steps", [
        "1. Immediate: Rotate the SECRET_KEY and use environment variables.",
        "2. Immediate: Disable Debug Mode.",
        "3. Short-term: Implement CSRF protection using Flask-WTF.",
        "4. Short-term: Pin dependency versions in requirements.txt.",
        "5. Long-term: Evaluate the need for user authentication."
    ])

    prs.save('Security_Review_Report.pptx')
    print("Presentation saved as Security_Review_Report.pptx")

if __name__ == "__main__":
    create_presentation()
